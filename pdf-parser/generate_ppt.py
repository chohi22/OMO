from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(5.625)

BG_COLOR = RGBColor(27, 42, 74)
TEXT_COLOR = RGBColor(255, 255, 255)
ACCENT_BLUE = RGBColor(74, 144, 217)
ACCENT_PURPLE = RGBColor(124, 58, 237)
ACCENT_AMBER = RGBColor(245, 158, 11)

def add_slide(layout_type=6):
    slide = prs.slides.add_slide(prs.slide_layouts[layout_type])
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = BG_COLOR
    return slide

def add_title(slide, text, top=0.3):
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(top), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.text = text
    title_frame.paragraphs[0].font.size = Pt(32)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = TEXT_COLOR
    
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0.5), Inches(top + 0.7), Inches(9), Inches(0.03)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = ACCENT_BLUE
    line.line.fill.background()

def add_bullet_text(text_frame, text, level=0, size=14, bold=False):
    p = text_frame.add_paragraph()
    p.text = text
    p.level = level
    p.font.size = Pt(size)
    p.font.color.rgb = TEXT_COLOR
    p.font.bold = bold
    return p

slide1 = add_slide()
title_box = slide1.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8), Inches(1))
tf = title_box.text_frame
tf.text = "PDF 텍스트 추출 시스템"
tf.paragraphs[0].font.size = Pt(48)
tf.paragraphs[0].font.bold = True
tf.paragraphs[0].font.color.rgb = TEXT_COLOR
tf.paragraphs[0].alignment = PP_ALIGN.CENTER

subtitle_box = slide1.shapes.add_textbox(Inches(1), Inches(2.8), Inches(8), Inches(0.6))
stf = subtitle_box.text_frame
stf.text = "PDF Parser and Web Viewer with PostgreSQL"
stf.paragraphs[0].font.size = Pt(20)
stf.paragraphs[0].font.color.rgb = ACCENT_BLUE
stf.paragraphs[0].alignment = PP_ALIGN.CENTER

date_box = slide1.shapes.add_textbox(Inches(1), Inches(4.5), Inches(8), Inches(0.4))
dtf = date_box.text_frame
dtf.text = "2026.02.13"
dtf.paragraphs[0].font.size = Pt(16)
dtf.paragraphs[0].font.color.rgb = RGBColor(200, 200, 200)
dtf.paragraphs[0].alignment = PP_ALIGN.CENTER

slide2 = add_slide()
add_title(slide2, "목차")
content_box = slide2.shapes.add_textbox(Inches(1.5), Inches(1.3), Inches(7), Inches(3.5))
tf = content_box.text_frame
tf.word_wrap = True

items = [
    "프로젝트 개요",
    "시스템 아키텍처",
    "주요 기능",
    "기술 스택",
    "DB 설계",
    "화면 구성",
    "데이터 흐름",
    "향후 계획"
]

for idx, item in enumerate(items):
    if idx == 0:
        tf.text = f"{idx+1}. {item}"
        tf.paragraphs[0].font.size = Pt(18)
        tf.paragraphs[0].font.color.rgb = TEXT_COLOR
        tf.paragraphs[0].space_after = Pt(16)
    else:
        p = tf.add_paragraph()
        p.text = f"{idx+1}. {item}"
        p.font.size = Pt(18)
        p.font.color.rgb = TEXT_COLOR
        p.space_after = Pt(16)

slide3 = add_slide()
add_title(slide3, "프로젝트 개요")

purpose_box = slide3.shapes.add_textbox(Inches(0.7), Inches(1.3), Inches(8.5), Inches(0.8))
ptf = purpose_box.text_frame
ptf.text = "목적"
ptf.paragraphs[0].font.size = Pt(20)
ptf.paragraphs[0].font.bold = True
ptf.paragraphs[0].font.color.rgb = ACCENT_AMBER

purpose_desc = slide3.shapes.add_textbox(Inches(0.7), Inches(1.7), Inches(8.5), Inches(0.6))
pdtf = purpose_desc.text_frame
pdtf.text = "PDF 문서에서 텍스트 추출 + 웹 조회/검색 시스템"
pdtf.paragraphs[0].font.size = Pt(16)
pdtf.paragraphs[0].font.color.rgb = TEXT_COLOR

features_box = slide3.shapes.add_textbox(Inches(0.7), Inches(2.5), Inches(8.5), Inches(0.5))
ftf = features_box.text_frame
ftf.text = "주요 특징"
ftf.paragraphs[0].font.size = Pt(20)
ftf.paragraphs[0].font.bold = True
ftf.paragraphs[0].font.color.rgb = ACCENT_AMBER

features_list = slide3.shapes.add_textbox(Inches(1.2), Inches(3.0), Inches(8), Inches(2))
fltf = features_list.text_frame
fltf.word_wrap = True

feature_items = [
    "PyMuPDF 고속 추출",
    "Flask 웹 서버",
    "PostgreSQL DB 영구 저장",
    "실시간 검색 하이라이트",
    "드래그앤드롭 업로드"
]

for idx, item in enumerate(feature_items):
    if idx == 0:
        fltf.text = f"• {item}"
        fltf.paragraphs[0].font.size = Pt(14)
        fltf.paragraphs[0].font.color.rgb = TEXT_COLOR
        fltf.paragraphs[0].space_after = Pt(10)
    else:
        p = fltf.add_paragraph()
        p.text = f"• {item}"
        p.font.size = Pt(14)
        p.font.color.rgb = TEXT_COLOR
        p.space_after = Pt(10)

slide4 = add_slide()
add_title(slide4, "시스템 아키텍처")

browser = slide4.shapes.add_shape(
    MSO_SHAPE.ROUNDED_RECTANGLE,
    Inches(1), Inches(1.5), Inches(1.5), Inches(0.8)
)
browser.fill.solid()
browser.fill.fore_color.rgb = ACCENT_PURPLE
browser.line.color.rgb = RGBColor(255, 255, 255)
browser.text_frame.text = "Browser"
browser.text_frame.paragraphs[0].font.size = Pt(14)
browser.text_frame.paragraphs[0].font.bold = True
browser.text_frame.paragraphs[0].font.color.rgb = TEXT_COLOR
browser.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
browser.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

flask = slide4.shapes.add_shape(
    MSO_SHAPE.ROUNDED_RECTANGLE,
    Inches(3.5), Inches(1.5), Inches(1.5), Inches(0.8)
)
flask.fill.solid()
flask.fill.fore_color.rgb = ACCENT_BLUE
flask.line.color.rgb = RGBColor(255, 255, 255)
flask.text_frame.text = "Flask:5000"
flask.text_frame.paragraphs[0].font.size = Pt(14)
flask.text_frame.paragraphs[0].font.bold = True
flask.text_frame.paragraphs[0].font.color.rgb = TEXT_COLOR
flask.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
flask.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

pymupdf = slide4.shapes.add_shape(
    MSO_SHAPE.ROUNDED_RECTANGLE,
    Inches(6), Inches(1.5), Inches(1.5), Inches(0.8)
)
pymupdf.fill.solid()
pymupdf.fill.fore_color.rgb = RGBColor(46, 125, 50)
pymupdf.line.color.rgb = RGBColor(255, 255, 255)
pymupdf.text_frame.text = "PyMuPDF"
pymupdf.text_frame.paragraphs[0].font.size = Pt(14)
pymupdf.text_frame.paragraphs[0].font.bold = True
pymupdf.text_frame.paragraphs[0].font.color.rgb = TEXT_COLOR
pymupdf.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
pymupdf.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

pdf_files = slide4.shapes.add_shape(
    MSO_SHAPE.ROUNDED_RECTANGLE,
    Inches(8), Inches(1.5), Inches(1.5), Inches(0.8)
)
pdf_files.fill.solid()
pdf_files.fill.fore_color.rgb = RGBColor(211, 47, 47)
pdf_files.line.color.rgb = RGBColor(255, 255, 255)
pdf_files.text_frame.text = "PDF Files"
pdf_files.text_frame.paragraphs[0].font.size = Pt(14)
pdf_files.text_frame.paragraphs[0].font.bold = True
pdf_files.text_frame.paragraphs[0].font.color.rgb = TEXT_COLOR
pdf_files.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
pdf_files.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

postgres = slide4.shapes.add_shape(
    MSO_SHAPE.ROUNDED_RECTANGLE,
    Inches(3.5), Inches(3.5), Inches(1.8), Inches(0.8)
)
postgres.fill.solid()
postgres.fill.fore_color.rgb = RGBColor(48, 63, 159)
postgres.line.color.rgb = RGBColor(255, 255, 255)
postgres.text_frame.text = "PostgreSQL:5432"
postgres.text_frame.paragraphs[0].font.size = Pt(14)
postgres.text_frame.paragraphs[0].font.bold = True
postgres.text_frame.paragraphs[0].font.color.rgb = TEXT_COLOR
postgres.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
postgres.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

arrow1 = slide4.shapes.add_connector(1, Inches(2.5), Inches(1.9), Inches(3.5), Inches(1.9))
arrow1.line.color.rgb = RGBColor(255, 255, 255)
arrow1.line.width = Pt(2)

arrow2 = slide4.shapes.add_connector(1, Inches(5), Inches(1.9), Inches(6), Inches(1.9))
arrow2.line.color.rgb = RGBColor(255, 255, 255)
arrow2.line.width = Pt(2)

arrow3 = slide4.shapes.add_connector(1, Inches(7.5), Inches(1.9), Inches(8), Inches(1.9))
arrow3.line.color.rgb = RGBColor(255, 255, 255)
arrow3.line.width = Pt(2)

arrow4 = slide4.shapes.add_connector(1, Inches(4.25), Inches(2.3), Inches(4.25), Inches(3.5))
arrow4.line.color.rgb = RGBColor(255, 255, 255)
arrow4.line.width = Pt(2)

slide5 = add_slide()
add_title(slide5, "주요 기능")

features_grid = [
    ("PDF 텍스트 추출", "PyMuPDF로 빠르고 정확한\n텍스트 추출", ACCENT_BLUE),
    ("웹 뷰어", "브라우저에서 PDF 내용\n실시간 확인", ACCENT_PURPLE),
    ("DB 저장", "PostgreSQL에 영구 저장\n메타데이터 관리", ACCENT_AMBER),
    ("검색 기능", "페이지별 검색 및\n하이라이트 표시", RGBColor(46, 125, 50))
]

positions = [
    (1, 1.5), (5.5, 1.5),
    (1, 3.3), (5.5, 3.3)
]

for idx, ((title, desc, color), (x, y)) in enumerate(zip(features_grid, positions)):
    box = slide5.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(x), Inches(y), Inches(3.5), Inches(1.3)
    )
    box.fill.solid()
    box.fill.fore_color.rgb = color
    box.line.color.rgb = RGBColor(255, 255, 255)
    box.line.width = Pt(2)
    
    tf = box.text_frame
    tf.text = title
    tf.paragraphs[0].font.size = Pt(18)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = TEXT_COLOR
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    p = tf.add_paragraph()
    p.text = desc
    p.font.size = Pt(12)
    p.font.color.rgb = RGBColor(230, 230, 230)
    p.alignment = PP_ALIGN.CENTER
    p.space_before = Pt(10)

slide6 = add_slide()
add_title(slide6, "기술 스택")

rows = 8
cols = 3
left = Inches(1.5)
top = Inches(1.5)
width = Inches(7)
height = Inches(3.5)

table = slide6.shapes.add_table(rows, cols, left, top, width, height).table

table.columns[0].width = Inches(2)
table.columns[1].width = Inches(3)
table.columns[2].width = Inches(2)

data = [
    ["구분", "기술", "버전"],
    ["언어", "Python", "3.12"],
    ["PDF 처리", "PyMuPDF", "1.27.1"],
    ["웹 프레임워크", "Flask", "3.1.2"],
    ["데이터베이스", "PostgreSQL", "16"],
    ["DB 드라이버", "psycopg2", "2.9.11"],
    ["프론트엔드", "Tailwind CSS", "2.2.19"],
    ["컨테이너", "Docker", "Latest"]
]

for row_idx, row_data in enumerate(data):
    for col_idx, cell_text in enumerate(row_data):
        cell = table.rows[row_idx].cells[col_idx]
        cell.text = cell_text
        
        if row_idx == 0:
            cell.fill.solid()
            cell.fill.fore_color.rgb = ACCENT_BLUE
            cell.text_frame.paragraphs[0].font.bold = True
            cell.text_frame.paragraphs[0].font.size = Pt(14)
        else:
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(50, 65, 95)
            cell.text_frame.paragraphs[0].font.size = Pt(12)
        
        cell.text_frame.paragraphs[0].font.color.rgb = TEXT_COLOR
        cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE

slide7 = add_slide()
add_title(slide7, "DB 설계")

table1 = slide7.shapes.add_shape(
    MSO_SHAPE.ROUNDED_RECTANGLE,
    Inches(1), Inches(1.5), Inches(3.5), Inches(2.5)
)
table1.fill.solid()
table1.fill.fore_color.rgb = RGBColor(50, 65, 95)
table1.line.color.rgb = ACCENT_BLUE
table1.line.width = Pt(2)

tf1 = table1.text_frame
tf1.word_wrap = True
tf1.text = "pdf_documents"
tf1.paragraphs[0].font.size = Pt(16)
tf1.paragraphs[0].font.bold = True
tf1.paragraphs[0].font.color.rgb = ACCENT_AMBER
tf1.paragraphs[0].alignment = PP_ALIGN.CENTER

fields1 = [
    "id (PK)",
    "filename (UNIQUE)",
    "file_size",
    "total_pages",
    "created_at"
]

for field in fields1:
    p = tf1.add_paragraph()
    p.text = f"  • {field}"
    p.font.size = Pt(11)
    p.font.color.rgb = TEXT_COLOR
    p.space_after = Pt(6)

table2 = slide7.shapes.add_shape(
    MSO_SHAPE.ROUNDED_RECTANGLE,
    Inches(5.5), Inches(1.5), Inches(3.5), Inches(2.5)
)
table2.fill.solid()
table2.fill.fore_color.rgb = RGBColor(50, 65, 95)
table2.line.color.rgb = ACCENT_PURPLE
table2.line.width = Pt(2)

tf2 = table2.text_frame
tf2.word_wrap = True
tf2.text = "pdf_pages"
tf2.paragraphs[0].font.size = Pt(16)
tf2.paragraphs[0].font.bold = True
tf2.paragraphs[0].font.color.rgb = ACCENT_AMBER
tf2.paragraphs[0].alignment = PP_ALIGN.CENTER

fields2 = [
    "id (PK)",
    "document_id (FK)",
    "page_number",
    "content (TEXT)",
    "created_at"
]

for field in fields2:
    p = tf2.add_paragraph()
    p.text = f"  • {field}"
    p.font.size = Pt(11)
    p.font.color.rgb = TEXT_COLOR
    p.space_after = Pt(6)

relation_line = slide7.shapes.add_connector(1, Inches(4.5), Inches(2.7), Inches(5.5), Inches(2.7))
relation_line.line.color.rgb = ACCENT_AMBER
relation_line.line.width = Pt(3)

rel_label = slide7.shapes.add_textbox(Inches(4.3), Inches(2.3), Inches(1.5), Inches(0.3))
rel_label.text_frame.text = "1 : N"
rel_label.text_frame.paragraphs[0].font.size = Pt(14)
rel_label.text_frame.paragraphs[0].font.bold = True
rel_label.text_frame.paragraphs[0].font.color.rgb = ACCENT_AMBER
rel_label.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

slide8 = add_slide()
add_title(slide8, "화면 구성")

screenshots = [
    ("/data/omo/test/screenshots/01_main.png", "메인 화면", 1),
    ("/data/omo/test/screenshots/02_pdf_view.png", "PDF 뷰 화면", 4),
    ("/data/omo/test/screenshots/03_search.png", "검색 화면", 7)
]

for img_path, caption, left_pos in screenshots:
    try:
        pic = slide8.shapes.add_picture(
            img_path,
            Inches(left_pos), Inches(1.5),
            width=Inches(2.5)
        )
        
        caption_box = slide8.shapes.add_textbox(
            Inches(left_pos), Inches(4.3), Inches(2.5), Inches(0.3)
        )
        caption_box.text_frame.text = caption
        caption_box.text_frame.paragraphs[0].font.size = Pt(12)
        caption_box.text_frame.paragraphs[0].font.color.rgb = ACCENT_BLUE
        caption_box.text_frame.paragraphs[0].font.bold = True
        caption_box.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    except:
        placeholder = slide8.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(left_pos), Inches(1.5), Inches(2.5), Inches(2.5)
        )
        placeholder.fill.solid()
        placeholder.fill.fore_color.rgb = RGBColor(60, 60, 60)
        placeholder.text_frame.text = caption
        placeholder.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        placeholder.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        placeholder.text_frame.paragraphs[0].font.color.rgb = TEXT_COLOR

slide9 = add_slide()
add_title(slide9, "데이터 흐름")

steps = [
    "1. PDF 업로드/스캔",
    "2. PyMuPDF 텍스트 추출",
    "3. pdf_documents 메타 저장",
    "4. pdf_pages 텍스트 저장",
    "5. 웹 뷰어 DB 조회 표시",
    "6. 실시간 검색 필터링"
]

colors = [ACCENT_PURPLE, ACCENT_BLUE, RGBColor(46, 125, 50), 
          ACCENT_AMBER, RGBColor(211, 47, 47), RGBColor(156, 39, 176)]

positions_flow = [
    (1, 1.5), (4.25, 1.5), (7.5, 1.5),
    (1, 3.5), (4.25, 3.5), (7.5, 3.5)
]

for idx, (step, color, (x, y)) in enumerate(zip(steps, colors, positions_flow)):
    box = slide9.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(x), Inches(y), Inches(2.5), Inches(0.7)
    )
    box.fill.solid()
    box.fill.fore_color.rgb = color
    box.line.color.rgb = RGBColor(255, 255, 255)
    box.line.width = Pt(2)
    
    box.text_frame.text = step
    box.text_frame.paragraphs[0].font.size = Pt(12)
    box.text_frame.paragraphs[0].font.bold = True
    box.text_frame.paragraphs[0].font.color.rgb = TEXT_COLOR
    box.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    box.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

for i in range(2):
    arrow = slide9.shapes.add_connector(1, Inches(3.5), Inches(1.85 + i*2), Inches(4.25), Inches(1.85 + i*2))
    arrow.line.color.rgb = RGBColor(200, 200, 200)
    arrow.line.width = Pt(2)

for i in range(2):
    arrow = slide9.shapes.add_connector(1, Inches(6.75), Inches(1.85 + i*2), Inches(7.5), Inches(1.85 + i*2))
    arrow.line.color.rgb = RGBColor(200, 200, 200)
    arrow.line.width = Pt(2)

arrow_down1 = slide9.shapes.add_connector(1, Inches(8.75), Inches(2.2), Inches(8.75), Inches(3.5))
arrow_down1.line.color.rgb = RGBColor(200, 200, 200)
arrow_down1.line.width = Pt(2)

slide10 = add_slide()
add_title(slide10, "향후 계획")

plans_box = slide10.shapes.add_textbox(Inches(1.5), Inches(1.5), Inches(7), Inches(3.5))
ptf = plans_box.text_frame
ptf.word_wrap = True

plans = [
    "OCR 이미지 PDF 지원",
    "테이블 구조 자동 인식",
    "메타데이터 추출 확장",
    "Full-Text Search 고도화",
    "REST API 추가",
    "사용자 인증"
]

for idx, plan in enumerate(plans):
    if idx == 0:
        ptf.text = f"• {plan}"
        ptf.paragraphs[0].font.size = Pt(18)
        ptf.paragraphs[0].font.color.rgb = TEXT_COLOR
        ptf.paragraphs[0].space_after = Pt(18)
    else:
        p = ptf.add_paragraph()
        p.text = f"• {plan}"
        p.font.size = Pt(18)
        p.font.color.rgb = TEXT_COLOR
        p.space_after = Pt(18)

prs.save('/data/omo/test/PDF_Parser_산출물.pptx')
print("✓ Presentation generated: PDF_Parser_산출물.pptx")
